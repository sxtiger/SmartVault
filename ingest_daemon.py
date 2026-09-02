#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
SmartVault / 智能仓库 —— 模块 A：自动化摄入与归档守护进程 (Ingest & Archive Daemon)
================================================================================
职责：
  1. 读取外部 config.json，注册多个 Obsidian Vault，用 watchdog 同时监听每个
     Vault 内部统一名称的「待处理笔记」收件箱目录。
  2. 实时捕获拖入的 Markdown 草稿及附件，按扩展名路由多模态解析：
       图像   (.png/.jpg/.jpeg/.heic...)   -> ocrmac（macOS Vision / Neural Engine）；
                                              auto 模式下 Vision 未识别到文字时 RapidOCR 兜底（手写）
       音视频 (.mp3/.m4a/.wav/.mp4/.mov..) -> mlx-whisper（Metal）/ openai-whisper（备选）
       PDF    (.pdf)                       -> PyMuPDF(fitz) 文本层；扫描页（中文手写）-> VLM(Qwen2.5-VL)
                                               逐字转写 + RapidOCR 兜底（v1.11.0）
       Office (.docx/.xlsx/.pptx)          -> python-docx / openpyxl / python-pptx
       iWork  (.pages/.numbers/.key)       -> 解剖 Zip 包/包目录 -> QuickLook/Preview.pdf -> fitz
  3. 动态扫描仓库目录树（一/二级）+ 读取 ai_context.md 人设与历史索引，注入 Prompt，
     呼叫本地 LM Studio 生成 Strict JSON（全字段纯简体中文约束）。
      草稿正文最前可放 ```smartvault 指令块（作者对本篇的归档要求），注入 Prompt 优先满足。
  4. 生成 YAML 属性 + 排版正文，移动附件与笔记至目标目录，清理草稿，
     反向追加 ai_context.md 历史索引，最后通过 obsidian:// URI 唤醒 Obsidian。

隐私承诺：全流程仅访问 localhost（LM Studio / 系统框架），零外部网络请求。

用法：
  python ingest_daemon.py --config config.json      # 前台常驻守护
  python ingest_daemon.py --check                   # 环境自检
  python ingest_daemon.py --scan                    # 一次性处理收件箱积压草稿后退出
  python ingest_daemon.py --once 草稿.md --vault 工作事务   # 调试单篇草稿
"""
from __future__ import annotations

import argparse
import base64
import json
import logging
import re
import shutil
import signal
import subprocess
import sys
import threading
import time
import zipfile
from dataclasses import dataclass
from datetime import datetime
from logging.handlers import RotatingFileHandler
from pathlib import Path
from queue import Empty, Queue
from typing import Any, Dict, List, Optional, Tuple
from urllib.parse import quote, unquote

# ---- 可选第三方依赖（容错导入，保证纯逻辑可在最小环境运行与自检）----
try:  # watchdog：仅常驻监听模式需要
    from watchdog.events import FileSystemEventHandler
    from watchdog.observers import Observer
except ImportError:  # pragma: no cover
    FileSystemEventHandler = object  # type: ignore[assignment,misc]
    Observer = None  # type: ignore[assignment]

try:  # requests：仅呼叫 LM Studio 时需要
    import requests
except ImportError:  # pragma: no cover
    requests = None  # type: ignore[assignment]

APP_NAME = "SmartVault"
LOG = logging.getLogger("smartvault.ingest")
DEFAULT_CONFIG = Path(__file__).resolve().parent / "config.json"

# ------------------------------------------------------------------ 扩展名路由表
IMAGE_EXTS = {".png", ".jpg", ".jpeg", ".heic", ".tiff", ".tif", ".bmp", ".gif", ".webp"}
AUDIO_EXTS = {".mp3", ".m4a", ".wav", ".flac", ".aac", ".ogg", ".wma", ".aiff"}
VIDEO_EXTS = {".mp4", ".mov", ".m4v", ".avi", ".mkv", ".webm"}
MEDIA_EXTS = AUDIO_EXTS | VIDEO_EXTS
IWORK_EXTS = {".pages", ".numbers", ".key"}
OFFICE_KIND_MAP = {".docx": "Word 文档", ".xlsx": "Excel 表格", ".pptx": "PPT 演示文稿"}
ATTACHMENT_EXTS = IMAGE_EXTS | MEDIA_EXTS | {".pdf"} | IWORK_EXTS | set(OFFICE_KIND_MAP) | {".txt"}

HIDDEN_DIR_NAMES = {".obsidian", ".trash", ".git", ".stfolder", ".smartvault"}

# LM Studio 结构化输出 JSON Schema（strict）
NOTE_JSON_SCHEMA: Dict[str, Any] = {
    "type": "object",
    "properties": {
        "target_folder": {"type": "string", "description": "目标目录，相对仓库根，如“一级目录”或“一级目录/二级目录”"},
        "new_filename": {"type": "string", "description": "新文件名，不含 .md 扩展名与路径分隔符"},
        "summary": {"type": "string", "description": "80~120 字摘要，必须严格取材于草稿原文，禁止出现原文没有的数字或事实"},
        "tags": {"type": "array", "items": {"type": "string"}, "minItems": 3, "maxItems": 5},
        "link_terms": {"type": "array", "items": {"type": "string"}, "maxItems": 8,
                       "description": "正文中实际出现的关键专有名词，逐字摘自正文，供系统确定性包裹双链"},
        "optimized_content": {"type": "string", "description": "排版后的完整 Markdown 正文"},
    },
    "required": ["target_folder", "new_filename", "summary", "tags", "link_terms", "optimized_content"],
    "additionalProperties": False,
}

# ------------------------------------------------------------------ 默认配置
CONFIG_DEFAULTS: Dict[str, Any] = {
    "lm_studio": {
        "base_url": "http://localhost:1234/v1",
        "chat_model": "qwen2.5-7b-instruct",
        "temperature": 0.3,
        "max_tokens": 4096,
        "timeout_seconds": 300,
        "structured_output": True,
    },
    "inbox_folder_name": "待处理笔记",
    "context_file": "ai_context.md",
    "ai_context_max_chars": 6000,
    "tree_depth": 2,
    "obsidian": {"wake_enabled": True},
    "vision": {"language_preference": ["zh-Hans", "en-US"]},
    "ocr": {"engine": "vlm", "image_engine": "auto", "pdf_dpi": 200,
            "pdf_max_ocr_pages": 20, "pdf_min_text_chars": 20, "rapidocr_max_width": 800,
            "vlm_model": "mlx-community/Qwen2.5-VL-7B-Instruct-4bit",
            "vlm_base_url": "http://localhost:1234/v1", "vlm_timeout": 300},
    "whisper": {"backend": "auto", "mlx_model": "mlx-community/whisper-large-v3-turbo",
                "openai_model": "small", "language": "zh"},
    "processing": {"debounce_seconds": 8, "quiet_seconds": 3, "attachment_wait_timeout": 30,
                   "attachments_subfolder": "附件", "allow_new_folder": True, "max_folder_depth": 2,
                   "fallback_folder": "未分类", "content_rewrite": False, "rewrite_max_chars": 6000,
                   "auto_wikilinks": True, "attachment_digest_enabled": True,
                   "attachment_digest_batch_chars": 4000},
    "limits": {"raw_note_max_chars": 30000, "attachment_max_chars": 100000,
               "attachment_prompt_max_chars": 12000},
    "rag": {"enabled": True, "embedding_model_path": "models/bge-small-zh-v1.5",
            "embedding_device": "mps", "chroma_dir": "data/chroma", "collection_name": "smartvault",
            "chunk_size": 500, "chunk_overlap": 80, "top_k": 4, "rescan_seconds": 300,
            "exclude_folders": [".obsidian", ".trash", "待处理笔记"]},
    "api": {"host": "127.0.0.1", "port": 8788},
    "vaults": [],
    "log_dir": "~/Library/Logs/SmartVault",
}


def _deep_merge(base: Dict[str, Any], override: Dict[str, Any]) -> Dict[str, Any]:
    """递归合并配置：override 覆盖 base（仅 dict 递归，list/标量直接替换）。"""
    out = dict(base)
    for k, v in (override or {}).items():
        if isinstance(v, dict) and isinstance(out.get(k), dict):
            out[k] = _deep_merge(out[k], v)
        else:
            out[k] = v
    return out


def load_config(path: Path) -> Dict[str, Any]:
    """读取 config.json 并与默认值合并；相对路径一律相对 config.json 所在目录解析。"""
    path = Path(path).expanduser().resolve()
    if not path.is_file():
        raise FileNotFoundError(f"配置文件不存在：{path}")
    with open(path, "r", encoding="utf-8") as f:
        user_cfg = json.load(f)
    cfg = _deep_merge(CONFIG_DEFAULTS, user_cfg or {})
    cfg["_config_dir"] = str(path.parent)
    # log_dir 支持 ~ 开头的绝对路径（v1.7.1 默认迁出 ~/Documents：TCC 禁止 launchd 打开其下日志）
    _log_dir = Path(cfg.get("log_dir", "~/Library/Logs/SmartVault")).expanduser()
    cfg["log_dir_abs"] = str((path.parent / _log_dir).resolve())
    return cfg


def setup_logging(cfg: Dict[str, Any], level: int = logging.INFO) -> None:
    log_dir = Path(cfg["log_dir_abs"])
    log_dir.mkdir(parents=True, exist_ok=True)
    fmt = logging.Formatter("%(asctime)s [%(levelname)s] %(name)s: %(message)s")
    LOG.setLevel(level)
    LOG.handlers.clear()
    fh = RotatingFileHandler(log_dir / "ingest_daemon.log", maxBytes=2 * 1024 * 1024,
                             backupCount=3, encoding="utf-8")
    fh.setFormatter(fmt)
    sh = logging.StreamHandler(sys.stdout)
    sh.setFormatter(fmt)
    LOG.addHandler(fh)
    LOG.addHandler(sh)
    LOG.propagate = False


# ================================================================== Vault 注册
@dataclass
class Vault:
    """一个被监听的 Obsidian 仓库。"""
    name: str          # Obsidian 中的仓库名（用于 URI 唤醒与展示）
    root: Path         # 仓库根目录
    inbox: Path        # 收件箱目录（待处理笔记）
    context_file: Path  # 仓库根下的 ai_context.md


def build_vaults(cfg: Dict[str, Any], strict: bool = False) -> List[Vault]:
    """根据 config.json 的 vaults 列表构建 Vault 对象；收件箱目录不存在则自动创建。"""
    inbox_name = cfg["inbox_folder_name"]
    ctx_name = cfg["context_file"]
    vaults: List[Vault] = []
    for item in cfg.get("vaults", []):
        root = Path(str(item["path"])).expanduser()
        name = str(item.get("name") or root.name)
        if not root.is_dir():
            msg = f"Vault 目录不存在，已跳过：{name} -> {root}"
            if strict:
                raise FileNotFoundError(msg)
            LOG.error(msg)
            continue
        inbox = root / inbox_name
        inbox.mkdir(parents=True, exist_ok=True)
        vaults.append(Vault(name=name, root=root, inbox=inbox, context_file=root / ctx_name))
    if not vaults:
        raise RuntimeError("config.json 中没有可用的 Vault（vaults 列表为空或路径全部无效）")
    return vaults


# ================================================================== 名称净化
_FILENAME_ILLEGAL_RE = re.compile(r'[\\/:*?"<>|#\^\[\]]')


def sanitize_component(text: str) -> str:
    """净化单个路径片段/文件名：去非法字符、折叠空白、去首尾空白与点号。"""
    text = _FILENAME_ILLEGAL_RE.sub("", str(text or ""))
    text = re.sub(r"\s+", " ", text).strip(" .")
    return text


def sanitize_filename(name: str, fallback: str = "未命名笔记", max_len: int = 80) -> str:
    """净化 LLM 生成的文件名（不含扩展名），保证跨平台安全。"""
    out = sanitize_component(name)
    return out[:max_len].strip(" .") or fallback


def sanitize_folder_parts(folder: str) -> List[str]:
    """把 LLM 生成的 target_folder 拆成安全的相对路径片段列表，拒绝绝对路径与穿越。"""
    parts = re.split(r"[/\\]+", str(folder or ""))
    parts = [sanitize_component(p) for p in parts]
    return [p for p in parts if p not in ("", ".", "..")]


def sanitize_tags(tags: List[str]) -> List[str]:
    out: List[str] = []
    for t in tags or []:
        t = re.sub(r"^#+", "", str(t).strip()).strip()
        t = re.sub(r"\s+", "-", t)
        if t and len(t) <= 24 and t not in out:
            out.append(t)
    return out[:5]


# ================================================================== 附件引用提取
# [[wikilink]] / ![[wikilink embed]]，兼容别名（|alias）与标题锚点（#section）
WIKILINK_RE = re.compile(r"\[\[([^\[\]]+?)\]\]")
# [text](path) / ![alt](path)，group(1)=方括号部分，group(2)=路径
MD_LINK_RE = re.compile(r"(\[[^\[\]]*\])\(\s*([^()\s]+)(?:\s+\"[^\"]*\")?\s*\)")
_URL_SCHEME_RE = re.compile(r"^[A-Za-z][A-Za-z0-9+.\-]*:")
# HTML 内嵌资源标签（<img>/<audio>/<video>/<source>/<embed> 的 src 属性），
# Kindle/HTML 转 Markdown 的产物常用此语法引用附件；双引号/单引号/裸值均兼容
HTML_SRC_RE = re.compile(
    r"""<(?:img|audio|video|source|embed)\b[^>]*?\bsrc\s*=\s*("[^"]*"|'[^']*'|[^\s>]+)""",
    re.IGNORECASE,
)


def find_attachment_refs(md_text: str) -> List[str]:
    """提取草稿中引用的附件文件名（按出现顺序去重）。

    - 支持 ![[文件.png]]、[[文件.pdf|别名]] 与 standard [文本](路径) / ![](路径) 语法
    - 支持 HTML 内嵌标签 <img src="路径">（Kindle/HTML 转 Markdown 产物常用）
    - 附件不限可解析类型：被链接即附件（.zip/.epub/.dwg 等任意扩展名均随笔记迁移），
      仅排除 .md 笔记链接、无扩展名双链（Obsidian 笔记链接）
    - 自动忽略：URL（http/obsidian:// 等带 scheme 的链接）、data: 内嵌资源
    """
    names: List[str] = []

    def _add(target: str) -> None:
        target = target.strip()
        if not target:
            return
        ext = Path(target).suffix.lower()
        if not ext or ext == ".md":  # 无扩展名双链与 .md 是笔记链接，不算附件
            return
        if target not in names:
            names.append(target)

    for m in WIKILINK_RE.finditer(md_text or ""):
        target = m.group(1).split("|", 1)[0].split("#", 1)[0].strip()
        _add(target)
    for m in MD_LINK_RE.finditer(md_text or ""):
        raw = m.group(2).strip()
        if _URL_SCHEME_RE.match(raw) or raw.startswith("#"):
            continue  # 跳过 http(s)、obsidian://、mailto: 等外部链接
        _add(unquote(Path(raw).name))
    for m in HTML_SRC_RE.finditer(md_text or ""):
        raw = m.group(1).strip("\"'").strip()
        if not raw or _URL_SCHEME_RE.match(raw) or raw.startswith("#"):
            continue  # 跳过 http(s)、data:image/... 内嵌资源等非收件箱附件
        _add(unquote(Path(raw).name))
    return names


def resolve_attachment(inbox: Path, ref_name: str) -> Optional[Path]:
    """在收件箱内定位附件文件：精确名 -> 无扩展名补全 -> 递归大小写不敏感匹配。"""
    ref_name = ref_name.strip()
    candidates = [ref_name]
    if not Path(ref_name).suffix:  # 用户写 ![[截图]] 但文件实际是 截图.png
        candidates += [Path(ref_name).stem + ext for ext in sorted(ATTACHMENT_EXTS)]
    for c in candidates:
        q = inbox / c
        if q.is_file():
            return q
    want = {c.lower() for c in candidates}
    for f in inbox.rglob("*"):
        try:
            if f.is_file() and not f.name.startswith(".") and f.name.lower() in want:
                return f
        except OSError:
            continue
    return None


# ================================================================== 多模态解析器
def _clip_text(text: str, limit: int) -> str:
    text = (text or "").strip()
    if len(text) <= limit:
        return text
    return text[:limit] + "\n……（内容超长，已截断）"


def _vision_image(path: Path, vision_cfg: Dict[str, Any]) -> str:
    """Vision OCR 原始文本（空串 = 未识别到文字）；占位说明由调用方统一生成。"""
    from ocrmac import ocrmac
    langs = vision_cfg.get("language_preference", ["zh-Hans", "en-US"])
    annotations = ocrmac.OCR(str(path), language_preference=langs).recognize()
    lines = []
    for item in annotations:
        lines.append(item[0] if isinstance(item, (tuple, list)) else str(item))
    return "\n".join(x.strip() for x in lines if x and x.strip())


def ocr_image(path: Path, vision_cfg: Dict[str, Any]) -> str:
    """图像 OCR：macOS Vision 框架（ocrmac），由 M4 Neural Engine 零延迟执行。"""
    return _vision_image(path, vision_cfg) or "（Vision OCR 未识别到文字，可能为纯图形图片）"


def _rapidocr_max_width(ocr_cfg: Optional[Dict[str, Any]]) -> int:
    """读 ocr.rapidocr_max_width（默认 800；<=0 = 关闭归一化；非法值回落 800）。"""
    try:
        v = int((ocr_cfg or {}).get("rapidocr_max_width", 800))
    except (TypeError, ValueError):
        return 800
    return v if v > 0 else 0


def _rapidocr_normalize(img: Any, max_width: int) -> Any:
    """RapidOCR 输入归一化（v1.10.0）：宽超 max_width 的图 LANCZOS 降采样后以 ndarray 传入。

    实测（真实潦草手写扫描页）：200dpi 渲染 1428px 宽直喂字准确率 52.5%，降采样至
    800px 后 67.6%（清晰字迹 97.1%→97.1% 无损）——PP-OCRv6 det/rec 对 ~600-1000px 宽
    输入最稳，大图直接喂反而丢细节。归一化遵循引擎自身行为：EXIF 方向转正与
    RGB→BGR（LoadImage 对路径/字节输入同样处理，ndarray 输入按 OpenCV 惯例）。
    max_width<=0 关闭（原样透传）；解码失败也透传原输入交引擎按原逻辑报错——绝不抛异常。
    """
    if max_width <= 0:
        return img
    try:
        import io
        import numpy as np
        from PIL import Image, ImageOps
        im = Image.open(io.BytesIO(img)) if isinstance(img, (bytes, bytearray)) else Image.open(img)
        im.load()
        if im.width <= max_width:
            return img
        im = ImageOps.exif_transpose(im) or im  # 手机照片 EXIF 方向（与引擎路径输入行为一致）
        scaled = im.resize((max_width, round(im.height * max_width / im.width)), Image.LANCZOS)
        return np.ascontiguousarray(np.asarray(scaled.convert("RGB"))[..., ::-1])  # RGB→BGR
    except Exception:  # noqa: BLE001 —— 归一化失败不拦路：透传原输入
        return img


def _rapidocr_image(path: Path, ocr_cfg: Optional[Dict[str, Any]] = None) -> str:
    """RapidOCR 识别整张图片（PP-OCRv6，中文手写友好）；返回原始文本（空串 = 无文字）。

    识别前按 ocr.rapidocr_max_width（默认 800）归一化输入宽度（见 _rapidocr_normalize）。
    """
    result = _get_rapidocr_engine()(
        _rapidocr_normalize(str(path), _rapidocr_max_width(ocr_cfg)))
    texts = [str(x).strip() for x in (getattr(result, "txts", None) or [])
             if x and str(x).strip()]
    return "\n".join(texts)


def extract_image(path: Path, cfg: Dict[str, Any]) -> Tuple[str, str]:
    """图像 OCR 引擎路由（v1.9.0）：ocr.image_engine = auto（默认）/ vision / rapidocr / off。

    auto：Vision 先识别（印刷体与规整手写效果好、Neural Engine 零延迟），无结果或异常时
    RapidOCR 兜底——潦草中文手写 Vision 可能整体识别失败，PP-OCRv6 更稳；纯图形图片两引擎
    都为空也只多花约 0.3s。任何失败降级为占位说明，不中断归档流水线。
    RapidOCR 路径（v1.10.0）：输入宽超 ocr.rapidocr_max_width（默认 800）自动 LANCZOS
    降采样——潦草手写字准确率实测 52.5%→67.6%，清晰字迹无损。
    注意：HEIC 仅 Vision 支持（OpenCV 不读 HEIC），HEIC 手写请保持 auto / vision 模式。
    """
    ocr_cfg = cfg.get("ocr") or {}
    engine = str(ocr_cfg.get("image_engine", "auto")).strip().lower()
    if engine in ("off", "disable", "false"):
        return "图像 OCR（已关闭）", "（图像 OCR 已关闭（ocr.image_engine=off），请打开原图片查看）"
    if engine == "vision":
        return "图像 OCR（Vision）", ocr_image(path, cfg.get("vision") or {})
    if engine == "rapidocr":
        try:
            text = _rapidocr_image(path, ocr_cfg)
        except RuntimeError as e:  # 引擎不可用 → 如实占位（不中断归档）
            return "图像 OCR（RapidOCR）", f"（{e}，请打开原图片查看）"
        return "图像 OCR（RapidOCR）", (text or "（RapidOCR 未识别到文字，可能为纯图形图片）")
    # auto：Vision 优先，失灵（无结果/异常）时 RapidOCR 兜底
    try:
        vision_text = _vision_image(path, cfg.get("vision") or {})
    except Exception as e:  # noqa: BLE001 —— Vision 失败不放弃，交 RapidOCR 救场
        LOG.warning("Vision OCR 失败（交 RapidOCR 兜底）：%s", e)
        vision_text = ""
    if vision_text:
        return "图像 OCR（Vision）", vision_text
    try:
        r_text = _rapidocr_image(path, ocr_cfg)
    except Exception as e:  # noqa: BLE001 —— 两引擎均失败：如实占位，不中断归档
        LOG.warning("RapidOCR 图像兜底失败：%s", e)
        return "图像 OCR（Vision→RapidOCR）", f"（Vision 未识别到文字，RapidOCR 兜底亦失败：{e}，请打开原图片查看）"
    if r_text:
        return "图像 OCR（Vision→RapidOCR 兜底）", r_text
    return "图像 OCR（Vision→RapidOCR）", "（Vision 与 RapidOCR 均未识别到文字，可能为纯图形图片）"


_OPENAI_WHISPER_MODEL: Dict[str, Any] = {}


def _transcribe_mlx(path: Path, wcfg: Dict[str, Any]) -> str:
    """MLX whisper：Apple Metal GPU 加速转写（M4 Pro 推荐）。"""
    import mlx_whisper
    result = mlx_whisper.transcribe(
        str(path),
        path_or_hf_repo=wcfg.get("mlx_model", "mlx-community/whisper-large-v3-turbo"),
        language=wcfg.get("language") or None,
    )
    return str(result.get("text", "")).strip()


def _transcribe_openai(path: Path, wcfg: Dict[str, Any]) -> str:
    """openai-whisper CPU 后端（需要 ffmpeg）。"""
    import whisper
    model_name = wcfg.get("openai_model", "small")
    if model_name not in _OPENAI_WHISPER_MODEL:
        LOG.info("加载 openai-whisper 模型：%s", model_name)
        _OPENAI_WHISPER_MODEL[model_name] = whisper.load_model(model_name, device="cpu")
    result = _OPENAI_WHISPER_MODEL[model_name].transcribe(
        str(path), language=wcfg.get("language") or None
    )
    return str(result.get("text", "")).strip()


def transcribe_media(path: Path, wcfg: Dict[str, Any]) -> str:
    """音视频转写：优先 MLX(Metal) 后端，失败自动回退 openai-whisper(CPU)。"""
    backend = str(wcfg.get("backend", "auto")).lower()
    errors: List[str] = []
    if backend in ("auto", "mlx"):
        try:
            return _transcribe_mlx(path, wcfg)
        except ImportError:
            if backend == "mlx":
                raise RuntimeError("未安装 mlx-whisper（pip install mlx-whisper）")
            errors.append("mlx-whisper 不可用")
        except Exception as e:  # noqa: BLE001 —— 单后端失败回退另一个，不中断流水线
            LOG.exception("MLX whisper 转写失败：%s", path.name)
            errors.append(f"mlx 异常: {e}")
    if backend in ("auto", "openai"):
        try:
            return _transcribe_openai(path, wcfg)
        except ImportError:
            errors.append("未安装 openai-whisper")
        except Exception as e:  # noqa: BLE001
            LOG.exception("openai-whisper 转写失败：%s", path.name)
            errors.append(f"openai 异常: {e}")
    raise RuntimeError("语音转写失败（" + "；".join(errors) + "），请安装 mlx-whisper 或 openai-whisper")


_RAPIDOCR_STATE: Dict[str, Any] = {}  # {"engine": 引擎单例|None, "error": 初始化失败原因|None}


def _get_rapidocr_engine():
    """RapidOCR 引擎惰性单例：模型内置于 wheel（PP-OCRv6，简体中文手写友好，离线可用）。

    首次构建约 1~2s，进程内复用；初始化失败（未安装/模型损坏）时缓存失败原因并
    以 RuntimeError 上抛，由调用方降级处理——OCR 永远不该中断归档流水线。
    """
    if "engine" not in _RAPIDOCR_STATE:
        try:
            from rapidocr import RapidOCR
            LOG.info("初始化 RapidOCR 引擎（PP-OCRv6，简体中文手写识别）")
            _RAPIDOCR_STATE["engine"] = RapidOCR()
        except ImportError:
            _RAPIDOCR_STATE["engine"] = None
            _RAPIDOCR_STATE["error"] = "未安装 rapidocr（pip install rapidocr）"
        except Exception as e:  # noqa: BLE001 —— 模型加载失败不中断流水线
            LOG.exception("RapidOCR 引擎初始化失败")
            _RAPIDOCR_STATE["engine"] = None
            _RAPIDOCR_STATE["error"] = str(e)
    if _RAPIDOCR_STATE["engine"] is None:
        raise RuntimeError(_RAPIDOCR_STATE.get("error") or "RapidOCR 初始化失败")
    return _RAPIDOCR_STATE["engine"]


_VLM_STATE: Dict[str, Any] = {}  # {"configured": 配置的 vlm_model, "resolved": 解析后的请求模型 id}

_VLM_OCR_PROMPT = (
    "逐字转写图片中的全部手写文字，要求：\n"
    "1. 从上到下、从左到右按行转写，不要遗漏任何一行（包括日期、落款、边角小字）\n"
    "2. 逐字照抄，看不清的字用「□」占位，绝不猜测编造\n"
    "3. 只输出转写文字本身，不要任何解释"
)


def _vlm_base_url(ocr_cfg: Optional[Dict[str, Any]] = None) -> str:
    """读 ocr.vlm_base_url（LM Studio OpenAI 兼容端点，默认本机 1234）。"""
    return str((ocr_cfg or {}).get("vlm_base_url",
                                   CONFIG_DEFAULTS["ocr"]["vlm_base_url"])).rstrip("/")


def _vlm_resolve_model(ocr_cfg: Optional[Dict[str, Any]] = None) -> str:
    """在 LM Studio /models 列表中解析 vlm_model 实际可用的请求 id（进程内缓存）。

    LM Studio 对同一模型可能以「发布方/全路径」或「目录短名」两种 id 提供服务：
    配置写全路径而实际以短名加载时，按「去发布方前缀 + 去量化后缀」做等价匹配
    （mlx-community/Qwen2.5-VL-7B-Instruct-4bit ≡ qwen2.5-vl-7b-instruct）。
    列表里找不到（服务未启/未加载）时原样返回配置值——LM Studio 支持按全路径
    JIT 拉起模型，成败交由识别请求本身定夺；任何解析异常同样不拦路。
    """
    configured = str((ocr_cfg or {}).get("vlm_model",
                                         CONFIG_DEFAULTS["ocr"]["vlm_model"])).strip()
    if _VLM_STATE.get("configured") == configured:
        return str(_VLM_STATE.get("resolved") or configured)
    resolved = configured
    try:
        resp = requests.get(f"{_vlm_base_url(ocr_cfg)}/models", timeout=5)
        resp.raise_for_status()
        ids = [str(m.get("id", "")) for m in ((resp.json() or {}).get("data") or [])]

        def _core(mid: str) -> str:  # 去发布方前缀与量化后缀，压成纯小写字母数字
            tail = re.sub(r"[-_](?:4|5|6|8)bit$", "", mid.split("/")[-1], flags=re.I)
            return re.sub(r"[^a-z0-9]", "", tail.lower())

        want = _core(configured)
        for mid in ids:
            if want and _core(mid) == want:
                resolved = mid
                break
    except Exception:  # noqa: BLE001 —— 解析失败不拦路：按配置值直接请求
        pass
    _VLM_STATE["configured"] = configured
    _VLM_STATE["resolved"] = resolved
    return resolved


def _vlm_ocr_png(png_bytes: bytes, ocr_cfg: Optional[Dict[str, Any]] = None) -> str:
    """VLM 逐字转写整页手写图（LM Studio 视觉对话接口）；返回原始文本（空串 = 未识别）。

    与 PP-OCRv6 不同，VLM 没有 800px 输入甜点——直接喂 pdf_dpi 原分辨率渲染
    （实测分辨率越高越稳）；temperature=0 求确定性。任何失败（连接拒绝/超时/
    模型未加载/HTTP 非 2xx/响应结构异常）抛异常，由调用方降级 RapidOCR 兜底
    ——VLM 永远不该中断归档流水线。注意：不走 apply_thinking_switch（/no_think
    软开关是 Qwen3 文本模型的约定，掺进转写指令会污染输出）。
    """
    if requests is None:
        raise RuntimeError("缺少依赖 requests（pip install requests）")
    cfg = ocr_cfg or {}
    try:
        timeout = int(cfg.get("vlm_timeout", 300))
    except (TypeError, ValueError):
        timeout = 300
    b64 = base64.b64encode(png_bytes).decode("ascii")
    payload: Dict[str, Any] = {
        "model": _vlm_resolve_model(cfg),
        "temperature": 0,
        "max_tokens": 2048,
        "stream": False,
        "messages": [{
            "role": "user",
            "content": [
                {"type": "image_url", "image_url": {"url": f"data:image/png;base64,{b64}"}},
                {"type": "text", "text": _VLM_OCR_PROMPT},
            ],
        }],
    }
    resp = requests.post(f"{_vlm_base_url(cfg)}/chat/completions", json=payload,
                         timeout=timeout)
    if resp.status_code != 200:
        raise RuntimeError(f"LM Studio VLM 接口 HTTP {resp.status_code}：{resp.text[:200]}")
    try:
        content = (((resp.json() or {}).get("choices") or [{}])[0]
                   .get("message", {}).get("content"))
    except Exception as e:  # noqa: BLE001 —— 响应结构异常视为失败，交调用方兜底
        raise RuntimeError(f"VLM 响应结构异常：{e}") from e
    text = content if isinstance(content, str) else ""
    return re.sub(r"^```[a-zA-Z]*\s*|```\s*$", "", text.strip()).strip()  # 去偶发代码围栏


def extract_pdf(path: Path, ocr_cfg: Optional[Dict[str, Any]] = None) -> Tuple[str, str]:
    """PDF 提取（v1.8.0 双通道 / v1.11.0 VLM 主引擎）：文本层优先（PyMuPDF）；
    扫描页（含中文手写）ocr.engine=vlm（默认）时经 LM Studio 用 Qwen2.5-VL
    逐字转写，失败/空结果时 RapidOCR（PP-OCRv6 + 输入归一化）兜底。

    逐页判定：文本层字符数 >= pdf_min_text_chars 直接取文本（电子 PDF 零成本、逐字保真）；
    低于阈值视为扫描页（手写笔记 PDF 即此类）——以 pdf_dpi 渲染成 PNG：VLM 拿原分辨率
    （与 PP-OCRv6 不同，VLM 无 800px 甜点，实测分辨率越高越稳）；RapidOCR 兜底前按
    ocr.rapidocr_max_width（默认 800）归一化（见 _rapidocr_normalize——归一化后潦草
    手写字准确率实测 52.5%→67.6%）。engine=rapidocr 跳过 VLM 纯走 RapidOCR。
    返回 (类型标签, 全文)；OCR 关闭/不可用/单页失败一律降级为占位说明，绝不中断归档。
    """
    import fitz
    cfg = ocr_cfg or {}
    dpi = max(72, int(cfg.get("pdf_dpi", 200)))
    max_w = _rapidocr_max_width(cfg)
    max_ocr_pages = max(0, int(cfg.get("pdf_max_ocr_pages", 20)))
    min_chars = max(1, int(cfg.get("pdf_min_text_chars", 20)))
    engine = str(cfg.get("engine", "vlm")).strip().lower()
    engine_on = engine not in ("off", "disable", "false")
    use_vlm = engine == "vlm"
    pages: List[str] = []
    vlm_pages = 0   # VLM 成功转写的扫描页数
    ocr_pages = 0   # RapidOCR 成功识别的扫描页数（rapidocr 模式或 VLM 降级兜底）
    with fitz.open(str(path)) as doc:
        for i, page in enumerate(doc, 1):
            t = page.get_text("text").strip()
            if len(t) >= min_chars:
                pages.append(f"[第 {i} 页]\n{t}")
                continue
            # —— 扫描页（无/极少文本层）：VLM 先转写，失败/空结果 RapidOCR 兜底
            if not engine_on:
                pages.append(f"[第 {i} 页]\n（扫描页：OCR 已关闭（ocr.engine=off），请打开原文件查看）")
                continue
            if vlm_pages + ocr_pages >= max_ocr_pages:
                pages.append(f"[第 {i} 页]\n（扫描页：超出 OCR 页数上限（{max_ocr_pages}），已省略，可调大 ocr.pdf_max_ocr_pages 后重试）")
                continue
            try:
                png_bytes = page.get_pixmap(dpi=dpi).tobytes("png")
            except Exception as e:  # noqa: BLE001 —— 单页渲染失败不影响其余页与流水线
                LOG.exception("PDF 页面渲染失败：%s 第 %d 页", path.name, i)
                pages.append(f"[第 {i} 页]\n（第 {i} 页渲染失败：{e}）")
                continue
            page_text = ""
            if use_vlm:
                try:
                    page_text = _vlm_ocr_png(png_bytes, cfg)
                    if page_text:
                        vlm_pages += 1
                except Exception as e:  # noqa: BLE001 —— VLM 失败降级 RapidOCR，不中断流水线
                    LOG.warning("VLM OCR 失败（%s 第 %d 页），降级 RapidOCR 兜底：%s",
                                path.name, i, e)
            if not page_text:
                try:
                    result = _get_rapidocr_engine()(_rapidocr_normalize(png_bytes, max_w))
                    texts = [str(x).strip() for x in (getattr(result, "txts", None) or [])
                             if x and str(x).strip()]
                    page_text = "\n".join(texts) or "（本页未识别到文字，可能为纯图形）"
                    ocr_pages += 1
                except RuntimeError as e:
                    pages.append(f"[第 {i} 页]\n（扫描页：{e}，请打开原文件查看）")
                    continue
                except Exception as e:  # noqa: BLE001 —— 单页 OCR 失败不影响其余页与流水线
                    LOG.exception("RapidOCR 识别失败：%s 第 %d 页", path.name, i)
                    pages.append(f"[第 {i} 页]\n（第 {i} 页 OCR 失败：{e}）")
                    continue
            pages.append(f"[第 {i} 页｜手写OCR]\n" + page_text)
    if use_vlm and (vlm_pages or ocr_pages):
        if ocr_pages:  # VLM 试过且有页靠 RapidOCR 兜底（失败/空结果降级）
            kind = f"PDF 提取（PyMuPDF+VLM→RapidOCR，手写 {vlm_pages + ocr_pages} 页）"
        else:
            kind = f"PDF 提取（PyMuPDF+VLM，手写 {vlm_pages} 页）"
    elif ocr_pages:
        kind = f"PDF 提取（PyMuPDF+RapidOCR，手写 {ocr_pages} 页）"
    else:
        kind = "PDF 文本（PyMuPDF）"
    return kind, "\n\n".join(pages) or "（PDF 中未提取到文本）"


def extract_docx(path: Path) -> str:
    from docx import Document
    doc = Document(str(path))
    parts: List[str] = [p.text.strip() for p in doc.paragraphs if p.text.strip()]
    for table in doc.tables:
        for row in table.rows:
            cells = [c.text.strip() for c in row.cells if c.text.strip()]
            if cells:
                parts.append(" | ".join(cells))
    return "\n".join(parts) or "（Word 文档为空）"


def extract_xlsx(path: Path) -> str:
    import openpyxl
    wb = openpyxl.load_workbook(str(path), read_only=True, data_only=True)
    parts: List[str] = []
    for ws in wb.worksheets:
        parts.append(f"[工作表：{ws.title}]")
        for row in ws.iter_rows(values_only=True):
            cells = [str(c).strip() for c in row if c is not None and str(c).strip()]
            if cells:
                parts.append(" | ".join(cells))
    wb.close()
    return "\n".join(parts) or "（Excel 表格为空）"


def extract_pptx(path: Path) -> str:
    from pptx import Presentation
    prs = Presentation(str(path))
    parts: List[str] = []
    for idx, slide in enumerate(prs.slides, 1):
        chunks: List[str] = [f"[幻灯片 {idx}]"]
        for shape in slide.shapes:
            if getattr(shape, "has_text_frame", False):
                t = shape.text_frame.text.strip()
                if t:
                    chunks.append(t)
            if getattr(shape, "has_table", False):
                for row in shape.table.rows:
                    cells = [c.text.strip() for c in row.cells if c.text.strip()]
                    if cells:
                        chunks.append(" | ".join(cells))
        try:
            notes = slide.notes_slide.notes_text_frame.text.strip()
            if notes:
                chunks.append(f"[备注] {notes}")
        except Exception:  # noqa: BLE001 —— 无备注页属正常
            pass
        parts.append("\n".join(chunks))
    return "\n\n".join(parts) or "（PPT 为空）"


def extract_iwork(path: Path) -> str:
    """Apple iWork：解剖 Zip 包/包目录 -> QuickLook/Preview.pdf -> fitz 提取预览文本。"""
    import fitz
    preview: Optional[bytes] = None
    if path.is_dir():  # Finder 中的“文件”在文件系统层可能是包目录
        q = path / "QuickLook" / "Preview.pdf"
        if q.is_file():
            preview = q.read_bytes()
    else:
        try:
            with zipfile.ZipFile(str(path)) as z:
                if "QuickLook/Preview.pdf" in z.namelist():
                    preview = z.read("QuickLook/Preview.pdf")
        except zipfile.BadZipFile:
            preview = None
    if not preview:
        return "（iWork 文件内未找到 QuickLook/Preview.pdf，建议在 Pages/Numbers/Keynote 中导出 PDF 后再拖入）"
    with fitz.open(stream=preview, filetype="pdf") as doc:
        pages = [p.get_text("text").strip() for p in doc]
    text = "\n\n".join(t for t in pages if t)
    return text or "（Preview.pdf 中未提取到文本）"


def read_text_file(path: Path) -> str:
    return path.read_text(encoding="utf-8", errors="replace").strip()


def dispatch_attachment(path: Path, cfg: Dict[str, Any]) -> Tuple[str, str]:
    """按扩展名路由解析器，返回 (类型标签, 解析文本)；任何失败都不会中断流水线。"""
    ext = path.suffix.lower()
    try:
        started = time.time()
        if ext in IMAGE_EXTS:
            kind, text = extract_image(path, cfg)
        elif ext in MEDIA_EXTS:
            kind, text = "音视频转录（whisper）", transcribe_media(path, cfg["whisper"])
        elif ext == ".pdf":
            kind, text = extract_pdf(path, cfg.get("ocr") or {})
        elif ext == ".docx":
            kind, text = OFFICE_KIND_MAP[ext], extract_docx(path)
        elif ext == ".xlsx":
            kind, text = OFFICE_KIND_MAP[ext], extract_xlsx(path)
        elif ext == ".pptx":
            kind, text = OFFICE_KIND_MAP[ext], extract_pptx(path)
        elif ext in IWORK_EXTS:
            kind, text = "iWork 预览（QuickLook）", extract_iwork(path)
        elif ext == ".txt":
            kind, text = "纯文本附件", read_text_file(path)
        else:
            # v1.6.4：不可解析类型（.zip/.epub/.dwg 等）不读内容（避免二进制乱码
            # 灌入 LLM），如实随笔记迁移，转录处仅标注占位
            kind = "暂不支持解析"
            text = f"（{path.name} 为「{ext}」类型，暂不支持自动解析；文件已随笔记归档至附件目录，请打开原文件查看）"
        LOG.info("附件解析完成 [%s] %s（%.1fs，%d 字符）", kind, path.name,
                 time.time() - started, len(text))
        return kind, text
    except Exception as e:  # noqa: BLE001
        LOG.exception("附件解析失败：%s", path.name)
        return "解析失败", f"（附件 {path.name} 解析失败：{e}）"


# ================================================================== 长附件分批提炼（v1.12.0 Map-Reduce）
_PAGE_NUM_RE = re.compile(r"\[第 (\d+) 页")
_PAGE_SPLIT_RE = re.compile(r"(?m)^(?=\[第 \d+ 页)")

_DIGEST_PROMPT = (
    "提炼以下附件转录文本的要点（300 字以内）：\n"
    "1. 提取关键事实、数字、人名、结论，逐字取材于原文，严禁编造\n"
    "2. 保持原文事实忠实，不改变原意\n"
    "3. 只输出要点本身，不要任何解释或前后缀"
)


def _split_transcript_pages(text: str) -> List[str]:
    """按 [第 N 页] / [第 N 页｜手写OCR] 页标记确定性切分转录文本。

    extract_pdf 输出天然携带该标记；无标记时整段作为单页返回。
    """
    if not text:
        return []
    parts = _PAGE_SPLIT_RE.split(text)
    return [p.strip() for p in parts if p.strip()]


def _extract_page_num(page_text: str) -> int:
    """从页面文本 [第 N 页] 标记中提取页码；无标记返回 0。"""
    m = _PAGE_NUM_RE.search(page_text)
    return int(m.group(1)) if m else 0


def _batch_pages(pages: List[str], batch_chars: int) -> List[Tuple[int, int, str]]:
    """贪心合并页面至 batch_chars 字符/批；超宽单页独立成批。

    返回 [(起始页码, 终止页码, 拼接文本)]；无页码标记的页面页码为 0。
    """
    if not pages or batch_chars <= 0:
        return []
    batches: List[Tuple[int, int, str]] = []
    cur: List[str] = []
    cur_len = 0
    cur_start = 0
    for page in pages:
        pnum = _extract_page_num(page)
        if not cur:
            cur = [page]
            cur_len = len(page)
            cur_start = pnum
        elif cur_len + len(page) + 2 > batch_chars:
            batches.append((cur_start, _extract_page_num(cur[-1]), "\n\n".join(cur)))
            cur = [page]
            cur_len = len(page)
            cur_start = pnum
        else:
            cur.append(page)
            cur_len += len(page) + 2  # +2 为 \n\n 分隔
    if cur:
        batches.append((cur_start, _extract_page_num(cur[-1]), "\n\n".join(cur)))
    return batches


def digest_attachment_text(text: str, cfg: Dict[str, Any],
                           client: LLMClient) -> Optional[str]:
    """分批提炼长附件转录文本（Map-Reduce），返回拼接要点或 None（整体失败）。

    触发：文本长度 > attachment_prompt_max_chars 且 attachment_digest_enabled=true。
    切分：按 [第 N 页] 页标记确定性切分，贪心合并至 batch_chars 字符/批。
    Map：每批调 LLM（thinking=False, max_tokens=1024, temperature=0.3）提炼要点，
         每批要点带页码范围标注（助元数据 LLM 建立文档结构感）。
    Reduce：各批要点拼接；若仍超注入上限再做一轮合并提炼，失败则头裁剪。
    降级：单批失败 → 该批降级为头 300 字切片；整体失败 → None（调用方回退头裁剪）。
    """
    proc = cfg.get("processing") or {}
    if not bool(proc.get("attachment_digest_enabled", True)):
        return None
    batch_chars = int(proc.get("attachment_digest_batch_chars", 4000))
    prompt_cap = int(cfg["limits"]["attachment_prompt_max_chars"])

    pages = _split_transcript_pages(text)
    if not pages:
        return None
    batches = _batch_pages(pages, batch_chars)
    if not batches:
        return None

    digests: List[str] = []
    total = len(batches)
    for idx, (start, end, batch_text) in enumerate(batches, 1):
        label = f"第 {start}-{end} 页" if start != end else f"第 {start} 页"
        LOG.info("附件要点提炼进度：%d/%d（%s，%d 字符）", idx, total, label, len(batch_text))
        try:
            msg = [{"role": "user", "content": f"{_DIGEST_PROMPT}\n\n---\n{batch_text}"}]
            result = client.chat(msg, max_tokens=1024, temperature=0.3, thinking=False)
            result = result.strip()
            if not result:
                LOG.warning("附件要点提炼批次 %s 返回空，降级头切片", label)
                result = _clip_text(batch_text, 300)
        except Exception as e:  # noqa: BLE001 —— 单批失败不中断整体
            LOG.warning("附件要点提炼批次 %s 失败，降级头切片：%s", label, e)
            result = _clip_text(batch_text, 300)
        digests.append(f"{label}：{result}")

    joined = "\n".join(digests)

    # Reduce：拼接后仍超注入上限（超大文档）→ 再做一轮合并提炼
    if len(joined) > prompt_cap:
        LOG.info("附件要点拼接 %d 字符仍超注入上限 %d，启动合并提炼", len(joined), prompt_cap)
        try:
            msg = [{"role": "user", "content": f"{_DIGEST_PROMPT}\n\n---\n{joined}"}]
            joined = client.chat(msg, max_tokens=1024, temperature=0.3,
                                 thinking=False).strip()
            if not joined:
                LOG.warning("附件要点合并提炼返回空，降级头裁剪")
                joined = _clip_text("\n".join(digests), prompt_cap)
        except Exception as e:  # noqa: BLE001 —— 合并失败不中断
            LOG.warning("附件要点合并提炼失败，降级头裁剪：%s", e)
            joined = _clip_text("\n".join(digests), prompt_cap)

    return joined if joined else None


def prepare_attachment_blocks(name: str, kind: str, text: str,
                              cfg: Dict[str, Any],
                              client: LLMClient) -> Tuple[List[str], str]:
    """准备附件的保全块列表与注入块（v1.12.0）。

    - 短附件（≤ attachment_prompt_max_chars）：直注原文，保全块 = [全文块]；
    - 长附件且提炼开启：提炼成功 → 保全块 = [速览块, 全文块]（速览在前、全文在后），
      注入块 = 要点提炼（原文 N 字，全文见文末转录块）；提炼失败 → 降级头裁剪；
    - 长附件且提炼关闭：保全块 = [全文块]，注入块 = 头裁剪。

    返回 (保全块列表, 注入块文本)；保全块列表按终稿落盘顺序排列。
    build_preserved_content 逐块渲染为独立 ``> [!quote]-`` 折叠块，无需改该函数。
    """
    limits = cfg["limits"]
    prompt_cap = int(limits["attachment_prompt_max_chars"])
    preserve_cap = int(limits["attachment_max_chars"])
    proc = cfg.get("processing") or {}
    digest_enabled = bool(proc.get("attachment_digest_enabled", True))

    # 短附件：直注原文
    if len(text) <= prompt_cap:
        block = f"◆ 附件「{name}」｜{kind}\n{text}"
        return [block], block

    # 长附件
    if digest_enabled:
        digest = digest_attachment_text(text, cfg, client)
        if digest:
            # 提炼成功：速览块在前、全文块在后（§2.3）
            speed_block = (f"◆ 附件「{name}」｜要点提炼"
                           f"（机器生成速览，若有出入以原附件为准）\n{digest}")
            full_block = f"◆ 附件「{name}」｜{kind}\n{_clip_text(text, preserve_cap)}"
            inject_block = (f"◆ 附件「{name}」｜要点提炼"
                            f"（原文 {len(text)} 字，全文见文末转录块）\n{digest}")
            LOG.info("附件要点速览块将随终稿落盘：%s（速览 %d 字 + 全文 %d 字）",
                     name, len(digest), min(len(text), preserve_cap))
            return [speed_block, full_block], inject_block
        LOG.warning("附件要点提炼整体失败，降级头裁剪注入：%s", name)

    # 长附件未提炼或提炼失败：头裁剪注入，全文保全
    full_block = f"◆ 附件「{name}」｜{kind}\n{_clip_text(text, preserve_cap)}"
    inject_block = f"◆ 附件「{name}」｜{kind}\n{_clip_text(text, prompt_cap)}"
    if len(text) > preserve_cap:
        LOG.info("附件保全上限触发：%s（%d 字 → %d 字）", name, len(text), preserve_cap)
    return [full_block], inject_block


# ================================================================== 目录树与上下文
def scan_tree(root: Path, depth: int = 2, exclude_names: frozenset = frozenset(),
              max_dirs: int = 200) -> str:
    """扫描仓库一/二级目录树，输出带树形符号的文本（供 LLM 选择归档位置）。"""
    lines: List[str] = [f"{root.name}/"]
    count = 0

    def rec(d: Path, level: int, prefix: str) -> None:
        nonlocal count
        if count >= max_dirs:
            return
        try:
            entries = sorted(
                (c for c in d.iterdir()
                 if c.is_dir() and not c.name.startswith(".")
                 and c.name not in HIDDEN_DIR_NAMES and c.name not in exclude_names),
                key=lambda x: x.name,
            )
        except OSError:
            return
        for i, c in enumerate(entries):
            if count >= max_dirs:
                lines.append(prefix + "……（目录过多已截断）")
                return
            last = i == len(entries) - 1
            lines.append(prefix + ("└── " if last else "├── ") + c.name + "/")
            count += 1
            if level < depth:
                rec(c, level + 1, prefix + ("    " if last else "│   "))

    rec(root, 1, "")
    return "\n".join(lines)


def load_ai_context(path: Path, max_chars: int) -> str:
    """读取 ai_context.md；超长时保留头部（规则区）+ 尾部（最近索引），中段省略。"""
    if not path.is_file():
        return ""
    try:
        text = path.read_text(encoding="utf-8", errors="replace")
    except OSError:
        return ""
    if len(text) <= max_chars:
        return text
    head_keep = max(1, int(max_chars * 0.35))
    tail_keep = max(1, max_chars - head_keep)
    return (text[:head_keep] + "\n……（中间历史索引已省略）……\n" + text[len(text) - tail_keep:])


def _drop_stale_ctx_entries(path: Path, stem: str) -> int:
    """移除 ai_context.md 中文件名（alias）与 stem 相同的「历史归档索引」条目。

    场景：误归档笔记移回收件箱重新归档时，指向旧目录的条目会随 Prompt
    注入形成「历史一致性」锚定（SYSTEM_PROMPT 规则 5），使 LLM 沿用旧
    目录、纠错失效；append-only 也会堆积重复条目。仅匹配标准生成行
    `- 文件：[[…|stem]]`（人工改写的非标准条目保守不动）；无移除不写盘；
    读写失败返回 0（追加逻辑照常执行）。
    """
    try:
        text = path.read_text(encoding="utf-8", errors="replace")
    except OSError:
        return 0
    parts = re.split(r"(?=^## )", text, flags=re.M)
    pat = re.compile(rf"^- 文件：\[\[.*\|{re.escape(stem)}\]\]\s*$", re.M)
    kept = [p for p in parts if not pat.search(p)]
    if len(kept) == len(parts):
        return 0
    try:
        path.write_text("".join(kept), encoding="utf-8")
    except OSError:
        return 0
    return len(parts) - len(kept)


def append_ai_context(vault: Vault, meta: Dict[str, Any], final_md: Path) -> None:
    """归档成功后向 ai_context.md 追加历史索引条目（文件不存在则带模板创建）。

    重归档自愈（v1.6.1）：追加前先移除同文件名的旧条目，索引始终只保留
    该笔记的最新归档位置。
    """
    path = vault.context_file
    rel = final_md.relative_to(vault.root)
    link = rel.as_posix()[:-3] if rel.as_posix().lower().endswith(".md") else rel.as_posix()
    tags_line = " ".join(f"#{t}" for t in meta.get("tags", []))
    entry = (
        f"\n## {datetime.now():%Y-%m-%d %H:%M}｜SmartVault 归档\n"
        f"- 文件：[[{link}|{final_md.stem}]]\n"
        f"- 目录：{meta.get('target_folder', '')}\n"
        f"- 摘要：{meta.get('summary', '')}\n"
        f"- 标签：{tags_line}\n"
    )
    if not path.exists():
        path.write_text(
            "# ai_context\n\n"
            "> 本文件由 SmartVault 维护：上半部为「AI 处理规则」（可人工编辑，"
            "守护进程每次都会注入给模型），下半部为「历史归档索引」（自动追加）。\n\n"
            "## AI 处理规则\n\n"
            "（示例）优先使用现有目录；标签使用小词而非长句；摘要聚焦事实。\n\n"
            "## 历史归档索引\n",
            encoding="utf-8",
        )
    else:
        removed = _drop_stale_ctx_entries(path, final_md.stem)
        if removed:
            LOG.info("ai_context 重归档去重：移除 %d 条旧条目（%s）", removed, final_md.stem)
    with open(path, "a", encoding="utf-8") as f:
        f.write(entry)


# ================================================================== LM Studio 客户端
def apply_thinking_switch(messages: List[Dict[str, str]], enabled: bool) -> List[Dict[str, str]]:
    """Qwen3 混合思考模型的软开关：enabled=False 时在最后一条 user 消息末尾追加 /no_think。

    Qwen3 官方 chat template 识别该标记后注入空 <think> 块，模型跳过思考直接作答
    （实测 qwen3-14b 同请求 16.2s -> 1.25s）。LM Studio 的 /v1 接口不支持
    chat_template_kwargs / enable_thinking 请求参数（实测均无效），软开关是唯一通道。
    返回新列表（浅拷贝逐条 dict），不修改调用方传入的消息。
    """
    if enabled:
        return messages
    out = [dict(m) for m in messages]
    for m in reversed(out):
        if m.get("role") == "user":
            m["content"] = f"{str(m.get('content') or '')}\n\n/no_think"
            break
    return out


class LLMClient:
    """LM Studio（OpenAI 兼容 /v1）客户端：优先 JSON Schema 结构化输出。"""

    def __init__(self, cfg: Dict[str, Any]):
        lm = cfg["lm_studio"]
        self.base_url = str(lm["base_url"]).rstrip("/")
        self.model = lm["chat_model"]
        # 采样参数对齐 Qwen3 官方 thinking 模式推荐值（temp 0.6 / top_p 0.95 /
        # top_k 20）；此前只发 temperature，top_p/top_k 落到 LM Studio 默认值
        # （1.0 / 40），属非官方组合。thinking=False 时经 /no_think 软开关跳过思考。
        self.temperature = float(lm.get("temperature", 0.6))
        self.top_p = float(lm.get("top_p", 0.95))
        self.top_k = int(lm.get("top_k", 20))
        self.thinking = bool(lm.get("thinking", True))
        self.max_tokens = int(lm.get("max_tokens", 4096))
        self.timeout = int(lm.get("timeout_seconds", 300))
        self.structured = bool(lm.get("structured_output", True))
        if requests is None:
            raise RuntimeError("缺少依赖 requests，请先 pip install requests")

    def chat(self, messages: List[Dict[str, str]],
             json_schema: Optional[Dict[str, Any]] = None,
             temperature: Optional[float] = None,
             max_tokens: Optional[int] = None,
             thinking: Optional[bool] = None) -> str:
        payload: Dict[str, Any] = {
            "model": self.model,
            "messages": apply_thinking_switch(
                messages, self.thinking if thinking is None else thinking),
            "temperature": self.temperature if temperature is None else temperature,
            "top_p": self.top_p,
            "top_k": self.top_k,
            "max_tokens": self.max_tokens if max_tokens is None else max_tokens,
            "stream": False,
        }
        if json_schema and self.structured:
            payload["response_format"] = {
                "type": "json_schema",
                "json_schema": {"name": "smartvault_note", "strict": True, "schema": json_schema},
            }
        url = f"{self.base_url}/chat/completions"
        resp = requests.post(url, json=payload, timeout=self.timeout)
        if resp.status_code == 400:
            body = resp.text[:400]
            try:
                err = resp.json().get("error")
                if isinstance(err, dict) and err.get("message"):
                    body = str(err["message"])
                elif err:
                    body = str(err)
            except Exception:  # noqa: BLE001
                pass
            if "exceeds the available context size" in body or (
                "context" in body.lower() and "tokens" in body.lower()
            ):
                # 上下文超限：与 response_format 无关，重发同样失败，直接给出可操作指引
                raise RuntimeError(
                    f"输入超过 LM Studio 上下文窗口（模型 {self.model}）：{body[:200]}。"
                    f"请以更大 context length 重载模型（如 lms load {self.model} -c 32768），"
                    "或拆分超长草稿；草稿已保留在收件箱。"
                )
            if "response_format" in payload:
                LOG.warning("LM Studio 不支持 response_format，回退纯提示词模式（原因：%s）", body[:120])
                payload.pop("response_format")
                resp = requests.post(url, json=payload, timeout=self.timeout)
        resp.raise_for_status()
        data = resp.json()
        return str(data["choices"][0]["message"]["content"] or "")

    def ping(self) -> bool:
        try:
            r = requests.get(f"{self.base_url}/models", timeout=5)
            return r.ok
        except Exception:  # noqa: BLE001
            return False


# ================================================================== Prompt 与 JSON 解析
SYSTEM_PROMPT = """你是「SmartVault 智能仓库」的资深知识管理图书管理员，负责把用户拖入收件箱的 Obsidian 草稿整理归档。

【输出契约（最高优先级）】
1. 你的回复必须是一个合法 JSON 对象本身：以 { 开始、以 } 结束；禁止 markdown 代码围栏、禁止任何解释性文字或前后缀。
2. JSON 字段固定为：
   - "target_folder": 字符串，笔记应归入的目标目录（相对仓库根目录，形如“一级目录”或“一级目录/二级目录”）。
   - "new_filename": 字符串，新文件名（不含 .md 扩展名，不含路径分隔符）。
   - "summary": 字符串，80~120 字的内容摘要；摘要中的事实与数字必须严格取材于草稿原文，严禁出现原文没有的数字、比例或结论。
   - "tags": 字符串数组，恰好 3~5 个主题标签，不带 # 号，不含空格。
   - "optimized_content": 字符串，整理排版后的完整 Markdown 正文；仅当系统未声明“原文保留模式”时才需要填写，声明后必须为空字符串 ""。
   - "link_terms": 字符串数组，正文中实际出现的关键专有名词（软件工具、通信协议、设备型号、技术名词、书名项目等），0~8 个、每个不超过 24 字符，必须逐字摘自草稿原文；系统将以确定性代码据此在正文上包裹 [[双链]]（不改任何其他字符）。

【整理规则】
1. target_folder 必须优先从用户给出的“仓库目录树”中选择已有目录；目录树中确无合适目录时，须依据笔记主题**新建简洁的一级目录**（2~6 个字，如“开发环境”“AI 工具”“网络工具”），让分类体系随归档自然生长，同主题笔记后续复用同一目录；最多二级深度；严禁使用“待处理笔记”、仓库根目录，以及“未分类”“笔记”“文档”“其他”等无信息量的目录名。
2. optimized_content 遵守 Obsidian Markdown 规范：文件内不要重复一级标题（标题由文件名承担），用二级/三级标题分节，善用列表与引用；正文中的所有事实、数字、百分比、指标、人名、结论必须逐字来自草稿原文或附件转录，严禁编造原文不存在的任何数字、比例或事实；整理仅限标题层级、列表化与删除冗余空白，禁止缩写、扩写或补充原文没有的内容；草稿为对话/问答体时必须保持原有问答结构与措辞，禁止重组为摘要式笔记。
3. 双链注入由系统代码完成：你只需在 link_terms 列出正文中实际出现的关键专有名词（软件工具、通信协议、设备型号、技术名词、书名项目，0~8 个，目录树已有目录名对应的实体优先）；严禁在 optimized_content 中自行添加、改写或删除 [[双链]]。
4. 附件的解析文本必须融入正文：以“> [!quote]- 附件：文件名”折叠引用块或独立小节呈现，冗长转录可提炼要点但不得丢失信息。
5. 若提供了 ai_context.md 内容，必须严格遵守其中的「AI 处理规则」，并与「历史归档索引」中已有标签体系、双链风格保持一致。
6. 全部输出内容（target_folder、new_filename、summary、tags、optimized_content）一律使用简体中文；专有名词、代码、命令、英文缩写、文件名与扩展名除外。
7. 清除草稿痕迹：删除“待处理”“测试”等临时字样与冗余空白，输出即终稿。
8. 超短草稿（正文不足约 200 字符，多为链接收藏、账号信息、碎片备忘）语义信号弱：必须依据笔记的实际用途与关键实体（链接指向的站点/工具、邮箱、账号、备忘主题）判断归类，不得凭个别词语的弱关联塞入已有目录，确无贴切目录时新建；这些关键信息须如实写入 summary 与 tags。
9. 若提供了【作者对本篇草稿的特别要求】，在不违反输出契约、事实忠实性与安全校验的前提下必须优先满足：作者指定的目标目录、文件名、标签与摘要侧重要求优先于 ai_context 规则与一般整理规则（作者指定的目录不存在时可按其指定名新建，仍受目录深度与禁用目录名约束）；确有无法满足或相互冲突的要求时按契约默认规则处理并在 summary 中如实反映，不得编造。"""

# ------------------------------------------------------------------ 草稿指令块（v1.9.0）
_FRONTMATTER_RE = re.compile(r"\A---[ \t]*\r?\n.*?\r?\n---[ \t]*(?:\r?\n|\Z)", re.DOTALL)
_INSTRUCTION_BLOCK_RE = re.compile(
    r"\A```smartvault[ \t]*\r?\n(.*?)\r?\n?```(?:[ \t]*\r?\n)*", re.DOTALL | re.IGNORECASE)
_REWRITE_DIRECTIVE_RE = re.compile(
    r"^[ \t]*(?:整理正文|重写正文|润色正文|content_rewrite)[ \t]*[:：=][ \t]*"
    r"(是|否|yes|no|true|false|on|off|开|关)[ \t]*$",
    re.IGNORECASE | re.MULTILINE)


def parse_draft_instructions(raw_md: str) -> Tuple[str, str]:
    """提取草稿指令块（v1.9.0）：正文最前面（frontmatter 之后、任何正文内容之前）的
    ```smartvault 围栏块，是作者对本篇草稿的归档要求（归入目录/文件名/标签等）。

    返回 (指令文本, 去除指令块后的正文)：
    - 指令块只认「正文开头」位置——正文中段出现的同名围栏块是普通内容，不提取、不动；
    - 未识别到（或内容为空）时正文原样返回（逐字不变，零副作用）；
    - 归档终稿不含指令块（此处确定性剥离，非 LLM 决定）；原稿完整备份于 .smartvault/backup/。
    """
    text = raw_md.lstrip("\ufeff")
    pos = 0
    fm = _FRONTMATTER_RE.match(text)
    if fm:
        pos = fm.end()
    rest = text[pos:]
    skip = re.match(r"(?:[ \t]*\r?\n)*", rest)   # frontmatter 与指令块之间的空行
    body_at = rest[skip.end():]
    m = _INSTRUCTION_BLOCK_RE.match(body_at)
    if m is None:
        return "", raw_md
    instructions = m.group(1).strip()
    if not instructions:
        return "", raw_md
    body = text[:pos] + body_at[m.end():]   # 指令块及其前后分隔空行一并剥离
    return instructions, body


def parse_rewrite_directive(instructions: str) -> Optional[bool]:
    """从草稿指令块解析「整理正文：是/否」结构化指令（唯一确定性覆盖项）。

    返回 True/False；指令块中未出现该行时返回 None（沿用全局 processing.content_rewrite）。
    正文是否允许 AI 整理是安全敏感项（v1.3.1 教训），自然语言请求不生效——
    必须用这一行结构化指令显式声明，且仍受 rewrite_max_chars 长度上限约束。
    """
    m = _REWRITE_DIRECTIVE_RE.search(instructions or "")
    if m is None:
        return None
    return m.group(1).strip().lower() in ("是", "yes", "true", "on", "开")


def build_user_prompt(vault_name: str, draft_name: str, raw_md: str,
                      attach_blocks: List[str], tree_text: str, ctx_text: str,
                      raw_max_chars: int, keep_original_content: bool = False,
                      instructions: str = "") -> str:
    parts = [
        f"【当前仓库：{vault_name}｜目录树（优先归类到已有目录）】",
        tree_text or "（目录树为空）",
        "【ai_context.md（仓库规则与历史索引）】",
        ctx_text or "（文件不存在，可自行判断）",
    ]
    if instructions:
        parts.append(
            "【作者对本篇草稿的特别要求（优先满足：优先级仅次于输出契约，高于 ai_context "
            "规则与一般整理规则；无法满足或与安全校验冲突时按契约默认规则处理）】\n" + instructions)
    parts.append(f"【草稿原文（原始文件名：{draft_name}）】")
    parts.append(_clip_text(raw_md, raw_max_chars))
    if attach_blocks:
        parts.append("【附件解析文本】\n" + "\n\n".join(attach_blocks))
    else:
        parts.append("【附件解析文本】\n（本篇草稿没有附件）")
    if keep_original_content:
        parts.append(
            "【本篇为原文保留模式（最高优先级，覆盖前述一切规则）】\n"
            f"本篇草稿共 {len(raw_md)} 字符，正文将由系统原样保留草稿原文，"
            "你绝对禁止对正文做任何摘要、压缩、改写或重排。因此：\n"
            "- optimized_content 字段必须返回空字符串 \"\"；\n"
            "- 你只需认真完成 target_folder、new_filename、summary、tags、link_terms 五个字段，"
            "其中 link_terms 供系统在保留的原文上以确定性代码包裹双链（不改任何其他字符）；"
            "其中 summary 也必须严格取材于原文，禁止出现原文没有的数字、比例或事实。"
        )
    parts.append("请依据以上信息，直接输出符合契约的 JSON 对象。")
    return "\n\n".join(parts)


def _strip_code_fence(text: str) -> str:
    text = text.strip()
    if text.startswith("```"):
        text = re.sub(r"^```[A-Za-z0-9_-]*[ \t]*\r?\n", "", text)
        text = re.sub(r"\r?\n?```\s*$", "", text)
    return text.strip()


def parse_llm_json(raw: str, fallback_filename: str = "未命名笔记") -> Dict[str, Any]:
    """鲁棒解析 LLM 输出为笔记元数据；容忍围栏、前后缀噪声与类型偏差。"""
    text = _strip_code_fence(raw or "")
    s, e = text.find("{"), text.rfind("}")
    if s < 0 or e <= s:
        raise ValueError(f"LLM 未返回 JSON 对象：{raw[:200]!r}")
    obj = json.loads(text[s:e + 1])

    folder = str(obj.get("target_folder") or "").strip()
    fname = str(obj.get("new_filename") or "").strip()
    summary = str(obj.get("summary") or "").strip()
    tags = obj.get("tags") or []
    if isinstance(tags, str):
        tags = [t for t in re.split(r"[,，;；\s]+", tags) if t]
    content = _strip_code_fence(str(obj.get("optimized_content") or ""))
    link_terms = obj.get("link_terms") or []
    if isinstance(link_terms, str):
        link_terms = [t for t in re.split(r"[,，;；\n]+", link_terms) if t.strip()]
    link_terms = [t.strip() for t in link_terms if str(t).strip()][:8]

    return {
        "target_folder": folder,
        "new_filename": fname or fallback_filename,
        "summary": summary or "（模型未生成摘要）",
        "tags": sanitize_tags([str(t) for t in tags]),
        "link_terms": link_terms,
        # 空 optimized_content 是合法值（长文保守模式），由 run_pipeline 回退为原文
        "optimized_content": content,
    }


# ================================================================== 落盘与归档
def _yaml_str(text: str) -> str:
    """YAML 双引号字符串安全转义（借用 JSON 转义规则）。"""
    return json.dumps(str(text), ensure_ascii=False)[1:-1]


def build_final_markdown(meta: Dict[str, Any], now: datetime) -> str:
    tags = meta.get("tags") or []
    tags_lines = "\n".join(f"  - {t}" for t in tags) or "  - 未分类"
    front = (
        "---\n"
        f"title: \"{_yaml_str(meta['new_filename'])}\"\n"
        f"date: {now:%Y-%m-%d %H:%M:%S}\n"
        f"summary: \"{_yaml_str(meta['summary'])}\"\n"
        "tags:\n"
        f"{tags_lines}\n"
        "source: SmartVault 自动归档\n"
        "---\n\n"
    )
    body = meta["optimized_content"].strip()
    if not body.endswith("\n"):
        body += "\n"
    return front + body


def choose_target_dir(vault: Vault, folder_str: str, proc: Dict[str, Any],
                      inbox_name: str) -> Path:
    """校验并确定目标目录：已有目录优先，允许受控新建，异常一律回退未分类。"""
    allow_new = bool(proc.get("allow_new_folder", True))
    max_depth = int(proc.get("max_folder_depth", 2))
    fallback = sanitize_component(proc.get("fallback_folder", "未分类")) or "未分类"
    parts = sanitize_folder_parts(folder_str)
    # LLM 常把仓库名一并输出（如“智能笔记/BMO/Profiles”），剥离后再判深度，避免误回退未分类
    while parts and parts[0] == vault.name:
        parts = parts[1:]
    if parts and inbox_name not in parts:
        cand = vault.root.joinpath(*parts)
        try:
            cand.relative_to(vault.root)
        except ValueError:
            cand = None  # 越界（不应发生，双保险）
        if cand is not None and cand != vault.root:
            if cand.is_dir():
                return cand
            if allow_new and len(parts) <= max_depth:
                cand.mkdir(parents=True, exist_ok=True)  # 深度已受限，可安全创建整条链路
                LOG.info("新建目标目录：%s", cand)
                return cand
    fb = vault.root / fallback
    fb.mkdir(parents=True, exist_ok=True)
    return fb


def prune_empty_dirs(inbox: Path) -> int:
    """清理收件箱内的空目录（自底向上，供归档后与启动补扫时调用）。

    归档移走附件后常残留空的 ``附件/`` 等目录。仅当目录为空、或只含
    Finder 元数据（.DS_Store）时才删除；收件箱根目录本身永不删除；
    含任何真实文件（含非 .DS_Store 的隐藏文件）的目录一律保留，
    因此不会误删正在等待附件的待处理草稿所在目录。返回删除的目录数。
    """
    try:
        dirs = [p for p in inbox.rglob("*") if p.is_dir()]
    except OSError:
        return 0
    removed = 0
    # 深度优先：先删空的子目录，父目录才可能随之变空
    for d in sorted(dirs, key=lambda p: len(p.parts), reverse=True):
        try:
            entries = list(d.iterdir())
        except OSError:
            continue
        for e in entries:
            if e.name != ".DS_Store":
                break
        else:
            for e in entries:
                try:
                    e.unlink()
                except OSError:
                    pass
            try:
                d.rmdir()
                removed += 1
            except OSError:
                pass
    return removed


def unique_path(p: Path) -> Path:
    """目标已存在时按 Obsidian 风格追加序号，避免覆盖。"""
    if not p.exists():
        return p
    for i in range(2, 1000):
        cand = p.with_name(f"{p.stem} {i}{p.suffix}")
        if not cand.exists():
            return cand
    return p.with_name(f"{p.stem} {int(time.time())}{p.suffix}")


def backup_draft(vault_root: Path, md_path: Path, keep: int = 100) -> Optional[Path]:
    """删除草稿前把原文备份到 vault 内 .smartvault/backup/，保留最近 keep 份。

    归档成功即删草稿是不可逆操作；备份兜底 LLM 整理失真或误归档的恢复需求。
    .smartvault 在 HIDDEN_DIR_NAMES 中（目录树扫描排除），也应配置进 RAG exclude_folders。
    """
    try:
        bdir = vault_root / ".smartvault" / "backup"
        bdir.mkdir(parents=True, exist_ok=True)
        stamp = datetime.now().strftime("%Y%m%d_%H%M%S")
        dest = unique_path(bdir / f"{stamp}_{md_path.name}")
        shutil.copy2(str(md_path), str(dest))
        olds = sorted([p for p in bdir.iterdir() if p.is_file()],
                      key=lambda p: p.name, reverse=True)
        for p in olds[keep:]:
            try:
                p.unlink()
            except OSError:
                pass
        return dest
    except Exception:  # noqa: BLE001 —— 备份失败不阻断归档，但必须留痕
        LOG.exception("草稿备份失败（继续归档）：%s", md_path)
        return None


def rewrite_links(content: str, moved_map: Dict[str, str], subfolder: str) -> str:
    """附件改名/移入子目录后，把正文中引用改写为新相对路径。

    - wikilink（[[x.png]]）Obsidian 全局按名解析，仅当附件被重命名时改写目标名；
    - standard link（[t](x.png)）与 HTML src（<img src="x.png">）是相对路径语法，
      统一改写为 子目录/新名（引号风格保持原样）。
    """
    if not moved_map:
        return content

    def rep_wikilink(m: "re.Match[str]") -> str:
        inner = m.group(1)
        target, _, alias = inner.partition("|")
        base, _, anchor = target.partition("#")
        if base.strip() in moved_map:
            new = moved_map[base.strip()]
            t = new + (f"#{anchor}" if anchor else "")
            return "[[" + t + (f"|{alias}" if alias else "") + "]]"
        return m.group(0)

    content = WIKILINK_RE.sub(rep_wikilink, content)

    def rep_mdlink(m: "re.Match[str]") -> str:
        raw = m.group(2)
        if _URL_SCHEME_RE.match(raw) or raw.startswith("#"):
            return m.group(0)
        name = Path(unquote(raw)).name
        if name in moved_map:
            new = moved_map[name]
            rel = f"{subfolder}/{new}" if subfolder else new
            return f"{m.group(1)}({rel})"
        return m.group(0)

    content = MD_LINK_RE.sub(rep_mdlink, content)

    def rep_htmlsrc(m: "re.Match[str]") -> str:
        quoted = m.group(1)
        q = quoted[0] if quoted[:1] in ('"', "'") else ""
        raw = quoted.strip("\"'").strip()
        if not raw or _URL_SCHEME_RE.match(raw) or raw.startswith("#"):
            return m.group(0)  # 外部 URL / data: 内嵌资源 / 锚点不动
        name = Path(unquote(raw)).name
        if name in moved_map:
            new = moved_map[name]
            rel = f"{subfolder}/{new}" if subfolder else new
            # 按 group(1) 区间精确重组，避免 alt="同值" 时误伤前序属性
            s, e = m.start(1) - m.start(0), m.end(1) - m.start(0)
            return m.group(0)[:s] + q + rel + q + m.group(0)[e:]
        return m.group(0)

    return HTML_SRC_RE.sub(rep_htmlsrc, content)


def wake_obsidian(vault_name: str, rel_path: str) -> None:
    """通过 obsidian:// URI 唤醒 Obsidian 并打开刚归档的笔记。"""
    uri = f"obsidian://open?vault={quote(vault_name, safe='')}&file={quote(rel_path, safe='')}"
    try:
        subprocess.Popen(["open", uri])
        LOG.info("已唤醒 Obsidian：%s", uri)
    except Exception as e:  # noqa: BLE001
        LOG.warning("唤醒 Obsidian 失败（不影响归档结果）：%s", e)


def _wait_file_stable(path: Path, checks: int = 10, interval: float = 1.0) -> bool:
    """等待附件文件大小稳定（大文件拖入时 Finder 是渐进拷贝）。"""
    last = -1
    for _ in range(checks):
        try:
            cur = path.stat().st_size
        except OSError:
            return False
        if cur == last and cur > 0:
            return True
        last = cur
        time.sleep(interval)
    return last > 0


# 双链注入的禁区：代码围栏、行内代码、md 链接 URL、HTML 标签（含 src 属性）
_WIKILINK_PROTECT_RE = re.compile(r"```.*?```|`[^`\n]*`|\]\([^)\s]*\)|<[^>\n]*>", re.S)


def apply_wikilinks(text: str, terms: List[str]) -> str:
    """确定性双链注入（v1.7.0）：把正文中出现的专有名词用 [[ ]] 包裹。

    LLM 只提供名词清单（link_terms），包裹由本函数以纯正则完成——除新增的
    [[ ]] 外不动任何字符（逐字保真可校验）；LLM 幻觉出正文没有的名词时
    正则匹配不到，自动忽略、零副作用。代码围栏/行内代码/链接 URL/HTML
    标签属禁区不注入；ASCII 词边界防止 RS485 误伤 RS4855；lookaround 防
    重复包裹已有 [[x]]；长词优先替换防止短词破坏长词。
    """
    if not text or not terms:
        return text
    cleaned: List[str] = []
    for t in terms:
        t = str(t).strip().strip("[]|").strip()
        if 1 <= len(t) <= 24 and "\n" not in t and t not in cleaned:
            cleaned.append(t)
    if not cleaned:
        return text
    cleaned.sort(key=len, reverse=True)          # 长词优先
    parts = _WIKILINK_PROTECT_RE.split(text)     # 禁区被剔除
    guards = _WIKILINK_PROTECT_RE.findall(text)  # 禁区原样放回
    rebuilt: List[str] = []
    for i, seg in enumerate(parts):
        for t in cleaned:
            seg = re.sub(rf"(?<![A-Za-z0-9_\[]){re.escape(t)}(?![A-Za-z0-9_\]])",
                         lambda _m, _t=t: f"[[{_t}]]", seg)
        rebuilt.append(seg)
        if i < len(guards):
            rebuilt.append(guards[i])
    return "".join(rebuilt)


def build_preserved_content(raw_md: str, attach_blocks: List[str]) -> str:
    """原文保留模式的正文组装：草稿原文逐字保留；附件转录以折叠引用块附加文末。

    附件转录是 OCR/Whisper 的机器产物而非用户原文：必须与原文明确区隔、可折叠、
    标注「以原附件为准」——既保证转录内容可被 RAG 检索，又避免转录误差污染原文。
    """
    if not attach_blocks:
        return raw_md
    lines = [raw_md.rstrip("\n"), "",
             "## 附：附件转录（机器自动生成，仅供检索，若有出入以原附件为准）", ""]
    for block in attach_blocks:
        header, _, text = block.partition("\n")
        lines.append(f"> [!quote]- {header}")
        body_lines = (text or "（空转录）").splitlines() or ["（空转录）"]
        lines.extend(f"> {ln}".rstrip() for ln in body_lines)
        lines.append("")
    return "\n".join(lines).rstrip("\n") + "\n"


# ================================================================== 归档流水线
def run_pipeline(cfg: Dict[str, Any], client: LLMClient, vault: Vault,
                 md_path: Path) -> Optional[Path]:
    """单篇草稿的完整处理链：读取 -> 等附件 -> 多模态解析 -> 上下文注入 ->
    LLM 提炼 -> 校验净化 -> 移动落盘 -> 更新索引 -> 唤醒 Obsidian。"""
    started = time.time()
    proc = cfg["processing"]
    inbox_name = cfg["inbox_folder_name"]
    raw = md_path.read_text(encoding="utf-8", errors="replace").lstrip("\ufeff")
    # v1.9.0 草稿指令块：正文最前的 ```smartvault 块是作者对本篇的归档要求——
    # 注入 Prompt 优先满足（目录/文件名/标签/摘要侧重等），指令块确定性剥离、不进终稿
    # （原稿完整备份于 .smartvault/backup/）；正文中间的同名围栏块是普通内容，不受影响
    instructions, raw = parse_draft_instructions(raw)
    if instructions:
        LOG.info("检测到草稿指令块（%d 字符，将作为作者特别要求注入 LLM，指令块不进入终稿）：\n%s",
                 len(instructions), instructions)
    LOG.info("开始处理草稿：%s（%d 字符）", md_path.name, len(raw))

    # 1) 等待引用的附件到齐（拖入多文件时附件可能晚于 md 到达）
    refs = find_attachment_refs(raw)
    deadline = time.time() + float(proc.get("attachment_wait_timeout", 30))
    while time.time() < deadline:
        missing = [r for r in refs if resolve_attachment(vault.inbox, r) is None]
        if not missing:
            break
        LOG.info("等待附件就绪：%s（剩余 %.0fs）", missing, deadline - time.time())
        time.sleep(2)

    # 2) 多模态解析附件（v1.12.0：保全与注入解耦 + 分批提炼）
    attach_blocks: List[str] = []        # 保全侧：文末转录折叠块（全文 + 要点速览）
    attach_blocks_prompt: List[str] = [] # 注入侧：LLM 元数据 Prompt（直注/提炼/头裁剪）
    attachments: List[Tuple[str, Path]] = []  # (引用名, 实际路径)
    for ref in refs:
        p = resolve_attachment(vault.inbox, ref)
        if p is None:
            LOG.warning("附件未找到，保留原引用：%s", ref)
            not_found = f"◆ 附件「{ref}」｜未找到\n（该附件未出现在收件箱中）"
            attach_blocks.append(not_found)
            attach_blocks_prompt.append(not_found)
            continue
        if not _wait_file_stable(p, checks=6, interval=1.0):
            LOG.warning("附件大小持续变化，按当前状态尽力解析：%s", p.name)
        kind, text = dispatch_attachment(p, cfg)
        attachments.append((ref, p))
        blocks, inject = prepare_attachment_blocks(p.name, kind, text, cfg, client)
        attach_blocks.extend(blocks)
        attach_blocks_prompt.append(inject)

    # 3) 组装动态上下文（原文保留模式：默认所有草稿正文逐字保留，LLM 只产元数据，
    #    杜绝任何改写导致的丢内容与幻觉——见 v1.3.1 事故复盘；短文 AI 润色为可选开关）
    tree_text = scan_tree(vault.root, depth=int(cfg.get("tree_depth", 2)),
                          exclude_names=frozenset({inbox_name}))
    ctx_text = load_ai_context(vault.context_file, int(cfg["ai_context_max_chars"]))
    rewrite_enabled = bool(proc.get("content_rewrite", False))   # 默认 False：正文永不改写
    per_note_rewrite = parse_rewrite_directive(instructions)
    if per_note_rewrite is not None:   # 指令块「整理正文：是/否」仅覆盖本篇
        LOG.info("草稿指令块「整理正文：%s」覆盖本篇正文改写开关（全局 content_rewrite=%s）",
                 "是" if per_note_rewrite else "否", rewrite_enabled)
        rewrite_enabled = per_note_rewrite
    rewrite_max = int(proc.get("rewrite_max_chars", 6000))
    keep_original = (not rewrite_enabled) or len(raw) > rewrite_max
    if per_note_rewrite and keep_original:
        LOG.warning("指令块要求整理正文，但正文 %d 字符超过 rewrite_max_chars（%d）上限，本篇仍保留原文",
                    len(raw), rewrite_max)
    user_prompt = build_user_prompt(vault.name, md_path.name, raw, attach_blocks_prompt,
                                    tree_text, ctx_text,
                                    int(cfg["limits"]["raw_note_max_chars"]),
                                    keep_original_content=keep_original,
                                    instructions=instructions)

    # 4) LLM 提炼（Strict JSON Schema 优先，失败自动回退纯提示词）
    raw_llm = client.chat([{"role": "system", "content": SYSTEM_PROMPT},
                           {"role": "user", "content": user_prompt}],
                          json_schema=NOTE_JSON_SCHEMA)
    meta = parse_llm_json(raw_llm, fallback_filename=md_path.stem)
    if keep_original:
        # 原文保留模式：正文逐字保留草稿原文（附件转录折叠附加于文末），无视 LLM 返回
        meta["optimized_content"] = build_preserved_content(raw, attach_blocks)
        LOG.info("原文保留模式（%d 字符，附件 %d 份）：正文保留原文，LLM 仅提供元数据",
                 len(raw), len(attach_blocks_prompt))
    else:
        # 短文模式：LLM 万一未返回正文也回退原文，任何情况下不允许内容丢失
        meta["optimized_content"] = meta["optimized_content"].strip() or raw
    LOG.info("LLM 提炼完成：目录=%s 文件名=%s 标签=%s", meta["target_folder"],
             meta["new_filename"], meta["tags"])

    # 5) 校验与净化
    folder = choose_target_dir(vault, meta["target_folder"], proc, inbox_name)
    now = datetime.now()
    final_md = unique_path(folder / (sanitize_filename(meta["new_filename"]) + ".md"))
    subfolder = str(proc.get("attachments_subfolder", "") or "").strip("/")

    # 6) 移动附件（含改名去重），并改写正文中的引用
    moved_map: Dict[str, str] = {}
    attach_dir = final_md.parent / subfolder if subfolder else final_md.parent
    if attachments:  # 无附件不预创建目录，避免遗留空 附件/ 目录
        attach_dir.mkdir(parents=True, exist_ok=True)
    for _, src in attachments:
        try:
            dest = unique_path(attach_dir / src.name)
            shutil.move(str(src), str(dest))
            moved_map[src.name] = dest.name
        except Exception:  # noqa: BLE001
            LOG.exception("附件移动失败（保留原处）：%s", src.name)
    meta["optimized_content"] = rewrite_links(meta["optimized_content"], moved_map, subfolder)

    # 6.5) 确定性双链注入（v1.7.0）：LLM 只提供名词表，包裹由代码完成，逐字保真
    if proc.get("auto_wikilinks", True) and meta.get("link_terms"):
        body = meta["optimized_content"]
        linked = apply_wikilinks(body, meta["link_terms"])
        if linked != body:
            LOG.info("双链注入：新增 %d 处 [[双链]]（确定性包裹，正文其余字符逐字不变）",
                     linked.count("[[") - body.count("[["))
        meta["optimized_content"] = linked

    # 7) 写入终稿、备份并清理草稿
    final_md.write_text(build_final_markdown(meta, now), encoding="utf-8")
    if (backup := backup_draft(vault.root, md_path)) is not None:
        LOG.info("草稿已备份：%s", backup.name)
    try:
        md_path.unlink()
    except OSError:
        LOG.exception("草稿清理失败（不影响归档）：%s", md_path)
    if (pruned := prune_empty_dirs(vault.inbox)) > 0:
        LOG.info("收件箱空目录清理：删除 %d 个（如归档后残留的空 附件/ 目录）", pruned)

    # 8) 反向写入 ai_context.md 历史索引 + 唤醒 Obsidian
    append_ai_context(vault, meta, final_md)
    rel = final_md.relative_to(vault.root).as_posix()
    if cfg.get("obsidian", {}).get("wake_enabled", True):
        wake_obsidian(vault.name, rel)

    LOG.info("归档完成：%s -> %s（耗时 %.1fs）", md_path.name, rel, time.time() - started)
    return final_md


# ================================================================== 监听与调度
class VaultState:
    """单个 Vault 的监听状态：防抖队列 + 处理工作线程。"""

    def __init__(self, cfg: Dict[str, Any], client: LLMClient, vault: Vault):
        self.cfg = cfg
        self.client = client
        self.vault = vault
        proc = cfg["processing"]
        self.debounce = float(proc.get("debounce_seconds", 8))
        self.quiet = float(proc.get("quiet_seconds", 3))
        self.pending: Dict[Path, Dict[str, Any]] = {}   # md -> {size, ready_at}
        self.lock = threading.Lock()
        self.last_event = time.time()
        self.queue: "Queue[Optional[Path]]" = Queue()
        self.stop_flag = threading.Event()
        self.worker = threading.Thread(target=self._worker, name=f"worker-{vault.name}", daemon=True)
        self.scheduler = threading.Thread(target=self._scheduler, name=f"sched-{vault.name}", daemon=True)

    # ---- 事件入口（由 watchdog handler 调用）----
    def on_event(self, path: Optional[Path], is_delete: bool) -> None:
        self.last_event = time.time()
        if path is None or path.name.startswith(".") or path.suffix.lower() != ".md":
            return
        with self.lock:
            if is_delete:
                self.pending.pop(path, None)
                return
            try:
                size = path.stat().st_size
            except OSError:
                size = 0
            self.pending[path] = {"size": size, "ready_at": time.time() + self.debounce}

    def seed_inbox(self) -> None:
        """启动时把收件箱中已存在的草稿纳入待处理（崩溃恢复）。"""
        for f in sorted(self.vault.inbox.rglob("*.md")):
            if f.name.startswith("."):
                continue
            self.on_event(f, is_delete=False)
        prune_empty_dirs(self.vault.inbox)

    # ---- 调度线程：等收件箱安静 + 文件大小稳定后入队 ----
    def _scheduler(self) -> None:
        while not self.stop_flag.is_set():
            time.sleep(1.0)
            now = time.time()
            quiet = now - self.last_event >= self.quiet
            with self.lock:
                items = list(self.pending.items())
            for md, rec in items:
                try:
                    if not md.exists():
                        with self.lock:
                            self.pending.pop(md, None)
                        continue
                    size = md.stat().st_size
                    if size != rec["size"]:  # 仍在写入，重置计时
                        rec["size"] = size
                        rec["ready_at"] = now + self.debounce
                        continue
                    if quiet and now >= rec["ready_at"]:
                        with self.lock:
                            self.pending.pop(md, None)
                        LOG.info("[%s] 草稿就绪，入队：%s", self.vault.name, md.name)
                        self.queue.put(md)
                except FileNotFoundError:
                    with self.lock:
                        self.pending.pop(md, None)
                except Exception:  # noqa: BLE001
                    LOG.exception("调度异常：%s", md)

    # ---- 工作线程：串行处理，避免同仓库并发写 ai_context.md ----
    def _worker(self) -> None:
        while not self.stop_flag.is_set():
            try:
                md = self.queue.get(timeout=1.0)
            except Empty:
                continue
            if md is None:
                break
            try:
                run_pipeline(self.cfg, self.client, self.vault, md)
            except Exception:  # noqa: BLE001 —— 失败保留草稿原处，便于人工重试
                LOG.exception("处理失败（草稿保留原处）：%s", md)
            finally:
                self.queue.task_done()

    def start(self) -> None:
        self.worker.start()
        self.scheduler.start()

    def stop(self) -> None:
        self.stop_flag.set()
        self.queue.put(None)


def make_event_handler(state: VaultState):
    """构造 watchdog 事件处理器：只关心收件箱内的 .md 增删改。"""
    if Observer is None:
        raise RuntimeError("缺少依赖 watchdog，请先 pip install watchdog")

    class InboxHandler(FileSystemEventHandler):  # type: ignore[misc,valid-type]
        def on_any_event(self, event) -> None:  # noqa: ANN001
            try:
                raw = getattr(event, "dest_path", None) or event.src_path
                if not raw:
                    return
                state.on_event(Path(raw), is_delete=(event.event_type == "deleted"))
            except Exception:  # noqa: BLE001
                LOG.exception("事件处理异常")

    return InboxHandler()


# ================================================================== 常驻守护
def run_daemon(cfg: Dict[str, Any]) -> None:
    vaults = build_vaults(cfg)
    client = LLMClient(cfg)
    if not client.ping():
        LOG.warning("LM Studio 未响应（%s），守护进程继续启动，处理时将自动重试",
                    cfg["lm_studio"]["base_url"])
    states = [VaultState(cfg, client, v) for v in vaults]
    observers = []
    for st in states:
        obs = Observer(timeout=1.0)
        obs.schedule(make_event_handler(st), str(st.vault.inbox), recursive=True)
        obs.start()
        observers.append(obs)
        st.seed_inbox()
        st.start()
        LOG.info("已监听 Vault [%s] 收件箱：%s", st.vault.name, st.vault.inbox)

    stop = threading.Event()

    def _sig(_sig_num, _frame) -> None:
        stop.set()

    signal.signal(signal.SIGINT, _sig)
    signal.signal(signal.SIGTERM, _sig)
    LOG.info("SmartVault 摄入守护进程运行中（监听 %d 个 Vault），Ctrl+C 退出", len(states))
    try:
        while not stop.is_set():
            time.sleep(0.5)
    finally:
        LOG.info("正在停止……")
        for st in states:
            st.stop()
        for obs in observers:
            obs.stop()
        for obs in observers:
            obs.join(timeout=5)
        for st in states:
            st.scheduler.join(timeout=5)
            st.worker.join(timeout=5)
        LOG.info("SmartVault 摄入守护进程已退出")


# ================================================================== 自检 / 批处理 / CLI
def run_check(cfg: Dict[str, Any]) -> int:
    """环境自检：逐项探测配置、依赖、LM Studio 与模型可用性。"""
    print(f"\n=== {APP_NAME} 环境自检 ===")
    ok_all = True

    def probe(label: str, fn) -> None:  # noqa: ANN001
        nonlocal ok_all
        try:
            msg = fn()
            print(f"  [✓] {label}" + (f"：{msg}" if msg else ""))
        except Exception as e:  # noqa: BLE001
            ok_all = False
            print(f"  [✗] {label}：{e}")

    probe("config.json 解析", lambda: f"{len(cfg.get('vaults', []))} 个 Vault")
    for item in cfg.get("vaults", []):
        p = Path(str(item["path"])).expanduser()
        probe(f"Vault 可达：{item.get('name', p.name)}",
              lambda p=p: (p.is_dir() or (_ for _ in ()).throw(FileNotFoundError(p))))
    probe("LM Studio 连通", lambda: (requests is not None and requests.get(
        cfg["lm_studio"]["base_url"].rstrip("/") + "/models", timeout=4).ok)
        or (_ for _ in ()).throw(RuntimeError(f"{cfg['lm_studio']['base_url']} 未响应")))
    if str((cfg.get("ocr") or {}).get("engine", "vlm")).strip().lower() == "vlm":
        def _vlm_probe() -> str:
            if requests is None:
                raise RuntimeError("缺少依赖 requests")
            r = requests.get(_vlm_base_url(cfg.get("ocr")) + "/models", timeout=4)
            r.raise_for_status()
            ids = [str(m.get("id", "")) for m in ((r.json() or {}).get("data") or [])]
            resolved = _vlm_resolve_model(cfg.get("ocr"))
            if resolved in ids:
                return f"模型已加载：{resolved}"
            return f"模型未预加载（{resolved}），首次识别将 JIT 拉起（约 30~60s）"
        probe("VLM 手写 OCR（LM Studio）", _vlm_probe)
    for mod, name in [("ocrmac", "图像 OCR（ocrmac）"), ("fitz", "PDF（PyMuPDF）"),
                      ("rapidocr", "手写 OCR 兜底（RapidOCR）"),
                      ("docx", "Word（python-docx）"), ("openpyxl", "Excel（openpyxl）"),
                      ("pptx", "PPT（python-pptx）")]:
        probe(name, lambda m=mod: __import__(m) and "")
    probe("whisper 后端", lambda: "")
    for backend in ("mlx_whisper", "whisper"):
        try:
            __import__(backend)
            print(f"      - 可用：{backend}")
        except ImportError:
            print(f"      - 未安装：{backend}")
    emb_abs = (Path(cfg["_config_dir"]) / cfg["rag"].get("embedding_model_path", "")).resolve()
    probe("本地嵌入模型目录", lambda: (emb_abs / "config.json").is_file()
          or (_ for _ in ()).throw(FileNotFoundError(f"{emb_abs}（需手动下载）")))
    print("=== 自检结束 ===\n")
    return 0 if ok_all else 1


def run_scan(cfg: Dict[str, Any]) -> int:
    """一次性处理所有收件箱中的积压草稿（处理完退出，适合手动批处理）。"""
    vaults = build_vaults(cfg)
    client = LLMClient(cfg)
    total = done = 0
    for vault in vaults:
        drafts = [f for f in sorted(vault.inbox.rglob("*.md")) if not f.name.startswith(".")]
        if not drafts:
            continue
        LOG.info("[%s] 扫描到 %d 篇积压草稿", vault.name, len(drafts))
        for f in drafts:
            total += 1
            try:
                run_pipeline(cfg, client, vault, f)
                done += 1
            except Exception:  # noqa: BLE001
                LOG.exception("处理失败（草稿保留原处）：%s", f)
    LOG.info("批处理完成：%d/%d 成功", done, total)
    return 0 if done == total else 1


def main() -> None:
    parser = argparse.ArgumentParser(description="SmartVault 摄入与归档守护进程")
    parser.add_argument("--config", default=str(DEFAULT_CONFIG), help="config.json 路径")
    parser.add_argument("--check", action="store_true", help="环境自检后退出")
    parser.add_argument("--scan", action="store_true", help="处理积压草稿后退出")
    parser.add_argument("--once", metavar="DRAFT_MD", help="调试：处理单篇草稿")
    parser.add_argument("--vault", metavar="NAME", help="配合 --once 指定 Vault 名称")
    parser.add_argument("--log-level", default="INFO")
    args = parser.parse_args()

    cfg = load_config(Path(args.config))
    setup_logging(cfg, getattr(logging, args.log_level.upper(), logging.INFO))
    if args.check:
        sys.exit(run_check(cfg))
    if args.scan:
        sys.exit(run_scan(cfg))
    if args.once:
        vaults = {v.name: v for v in build_vaults(cfg)}
        key = args.vault or ""
        if key not in vaults:
            raise SystemExit(f"--vault 必须是以下之一：{', '.join(vaults)}")
        draft = Path(args.once).expanduser()
        if not draft.is_absolute():
            draft = vaults[key].inbox / draft
        if not draft.is_file():
            raise SystemExit(f"草稿不存在：{draft}")
        sys.exit(0 if run_pipeline(cfg, LLMClient(cfg), vaults[key], draft) else 1)
    run_daemon(cfg)


if __name__ == "__main__":
    main()










